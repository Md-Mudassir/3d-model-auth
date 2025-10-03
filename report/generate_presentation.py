"""
Script to generate a professional PowerPoint presentation for the 3D Model Authentication project
Creates a 9-slide presentation covering all key aspects of the project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate a professional PowerPoint presentation about the 3D Model Authentication project"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(31, 78, 121)  # Dark blue
    ACCENT_COLOR = RGBColor(68, 114, 196)  # Medium blue
    TEXT_COLOR = RGBColor(64, 64, 64)  # Dark gray
    
    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "3D Model Digital Signature System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Using Steganographic Authentication & RSA Cryptography"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = ACCENT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add author info
    author_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
    author_frame = author_box.text_frame
    author_frame.text = "CHRIST (Deemed to be University)\nDepartment of Computer Science"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(16)
    author_para.font.color.rgb = TEXT_COLOR
    author_para.alignment = PP_ALIGN.CENTER
    
    # ==================== SLIDE 2: PROBLEM STATEMENT ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    title2 = slide2.shapes.title
    title2.text = "Problem Statement"
    title2.text_frame.paragraphs[0].font.size = Pt(40)
    title2.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "Challenges in 3D Content Protection"
    
    p2_1 = tf2.add_paragraph()
    p2_1.text = "Piracy & Unauthorized Distribution"
    p2_1.level = 1
    p2_1.font.size = Pt(20)
    
    p2_2 = tf2.add_paragraph()
    p2_2.text = "Loss of Artist Attribution"
    p2_2.level = 1
    p2_2.font.size = Pt(20)
    
    p2_3 = tf2.add_paragraph()
    p2_3.text = "Difficulty Detecting Tampering"
    p2_3.level = 1
    p2_3.font.size = Pt(20)
    
    p2_4 = tf2.add_paragraph()
    p2_4.text = "No Reliable Authentication Method"
    p2_4.level = 1
    p2_4.font.size = Pt(20)
    
    p2_5 = tf2.add_paragraph()
    p2_5.text = "Need for embedded, tamper-resistant signatures"
    p2_5.level = 1
    p2_5.font.size = Pt(20)
    p2_5.font.italic = True
    p2_5.font.color.rgb = ACCENT_COLOR
    
    # ==================== SLIDE 3: PROPOSED SOLUTION ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title3 = slide3.shapes.title
    title3.text = "Proposed Solution"
    title3.text_frame.paragraphs[0].font.size = Pt(40)
    title3.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Digital Signature + Steganography"
    tf3.paragraphs[0].font.size = Pt(24)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    p3_1 = tf3.add_paragraph()
    p3_1.text = "RSA Cryptography (2048-bit keys)"
    p3_1.level = 1
    p3_1.font.size = Pt(20)
    
    p3_2 = tf3.add_paragraph()
    p3_2.text = "Generate unique digital signatures for each artist"
    p3_2.level = 2
    p3_2.font.size = Pt(18)
    
    p3_3 = tf3.add_paragraph()
    p3_3.text = "Vertex-Level Steganography"
    p3_3.level = 1
    p3_3.font.size = Pt(20)
    
    p3_4 = tf3.add_paragraph()
    p3_4.text = "Embed signatures in 3D model geometry (LSB manipulation)"
    p3_4.level = 2
    p3_4.font.size = Pt(18)
    
    p3_5 = tf3.add_paragraph()
    p3_5.text = "Artist Attribution System"
    p3_5.level = 1
    p3_5.font.size = Pt(20)
    
    p3_6 = tf3.add_paragraph()
    p3_6.text = "Database-backed profiles with metadata embedding"
    p3_6.level = 2
    p3_6.font.size = Pt(18)
    
    # ==================== SLIDE 4: SYSTEM ARCHITECTURE ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title4 = slide4.shapes.title
    title4.text = "System Architecture"
    title4.text_frame.paragraphs[0].font.size = Pt(40)
    title4.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "Modular Design with 5 Core Components"
    tf4.paragraphs[0].font.size = Pt(22)
    tf4.paragraphs[0].font.bold = True
    
    components = [
        ("User Interface Layer", "Streamlit-based web interface"),
        ("Artist Management Module", "Profile creation & key management"),
        ("Cryptographic Module", "RSA signature generation & verification"),
        ("Steganographic Module", "Embedding & extraction algorithms"),
        ("Database Layer", "SQLite for persistent storage")
    ]
    
    for comp_name, comp_desc in components:
        p = tf4.add_paragraph()
        p.text = f"{comp_name}"
        p.level = 1
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        p_desc = tf4.add_paragraph()
        p_desc.text = comp_desc
        p_desc.level = 2
        p_desc.font.size = Pt(16)
    
    # ==================== SLIDE 5: KEY FEATURES ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title5 = slide5.shapes.title
    title5.text = "Key Features"
    title5.text_frame.paragraphs[0].font.size = Pt(40)
    title5.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    
    features = [
        ("🔐 Tamper-Resistant Signatures", "Distributed across multiple vertices"),
        ("👨‍🎨 Multi-Artist Support", "Unique cryptographic identity per artist"),
        ("🚫 Re-Signing Prevention", "Protects original artist attribution"),
        ("✅ Signature Verification", "Detect unauthorized modifications"),
        ("🎨 Visual Quality Preservation", "Imperceptible changes to geometry"),
        ("💾 Database-Backed", "Persistent artist profiles in SQLite"),
        ("🖥️ Interactive 3D Viewer", "Black background, white models, enhanced lighting")
    ]
    
    tf5.text = features[0][0]
    tf5.paragraphs[0].font.size = Pt(18)
    tf5.paragraphs[0].font.bold = True
    
    p_desc = tf5.add_paragraph()
    p_desc.text = features[0][1]
    p_desc.level = 1
    p_desc.font.size = Pt(16)
    
    for feat_title, feat_desc in features[1:]:
        p = tf5.add_paragraph()
        p.text = feat_title
        p.font.size = Pt(18)
        p.font.bold = True
        
        p_d = tf5.add_paragraph()
        p_d.text = feat_desc
        p_d.level = 1
        p_d.font.size = Pt(16)
    
    # ==================== SLIDE 6: TECHNICAL IMPLEMENTATION ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title6 = slide6.shapes.title
    title6.text = "Technical Implementation"
    title6.text_frame.paragraphs[0].font.size = Pt(40)
    title6.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Technology Stack"
    tf6.paragraphs[0].font.size = Pt(24)
    tf6.paragraphs[0].font.bold = True
    tf6.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    tech_stack = [
        ("Python 3.8+", "Core implementation language"),
        ("Streamlit", "Web application framework"),
        ("cryptography library", "RSA-PSS with SHA-256"),
        ("SQLite", "Artist profile database"),
        ("Three.js", "3D model visualization"),
        ("OBJ Format", "Widely-supported 3D file format")
    ]
    
    for tech, desc in tech_stack:
        p = tf6.add_paragraph()
        p.text = f"{tech}: {desc}"
        p.level = 1
        p.font.size = Pt(18)
    
    # ==================== SLIDE 7: STEGANOGRAPHIC ALGORITHM ====================
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title7 = slide7.shapes.title
    title7.text = "Steganographic Algorithm"
    title7.text_frame.paragraphs[0].font.size = Pt(40)
    title7.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "Embedding Process"
    tf7.paragraphs[0].font.size = Pt(24)
    tf7.paragraphs[0].font.bold = True
    tf7.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    embed_steps = [
        "Parse OBJ file and identify all vertices",
        "Convert signature to binary representation",
        "Split each byte into four 2-bit pairs",
        "Modify z-coordinates using LSB manipulation",
        "Embed artist metadata in Base64 format",
        "Add authentication marker to file header"
    ]
    
    for i, step in enumerate(embed_steps, 1):
        p = tf7.add_paragraph()
        p.text = f"{i}. {step}"
        p.level = 1
        p.font.size = Pt(18)
    
    # Add extraction note
    p_extract = tf7.add_paragraph()
    p_extract.text = "Extraction: Reverse process to reconstruct signature"
    p_extract.font.size = Pt(18)
    p_extract.font.italic = True
    p_extract.font.color.rgb = ACCENT_COLOR
    
    # ==================== SLIDE 8: SECURITY FEATURES ====================
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title8 = slide8.shapes.title
    title8.text = "Security Features"
    title8.text_frame.paragraphs[0].font.size = Pt(40)
    title8.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "Multi-Layer Protection"
    tf8.paragraphs[0].font.size = Pt(24)
    tf8.paragraphs[0].font.bold = True
    
    security_features = [
        ("Cryptographic Integrity", "RSA-PSS with 2048-bit keys, SHA-256 hashing"),
        ("Tamper Detection", "Signature verification detects any modifications"),
        ("Re-Signing Prevention", "Blocks attempts to override existing signatures"),
        ("Artist Verification", "Embedded metadata confirms original creator"),
        ("Duplicate Prevention", "No duplicate artist names in registry"),
        ("Distributed Storage", "Signature spread across multiple vertices")
    ]
    
    for sec_title, sec_desc in security_features:
        p = tf8.add_paragraph()
        p.text = sec_title
        p.level = 1
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        
        p_d = tf8.add_paragraph()
        p_d.text = sec_desc
        p_d.level = 2
        p_d.font.size = Pt(16)
    
    # ==================== SLIDE 9: CONCLUSION & FUTURE WORK ====================
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    
    title9 = slide9.shapes.title
    title9.text = "Conclusion & Future Work"
    title9.text_frame.paragraphs[0].font.size = Pt(40)
    title9.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "Achievements"
    tf9.paragraphs[0].font.size = Pt(22)
    tf9.paragraphs[0].font.bold = True
    tf9.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    achievements = [
        "✓ Robust authentication system for 3D models",
        "✓ Imperceptible steganographic embedding",
        "✓ User-friendly interface for creators",
        "✓ Strong tamper-resistance mechanisms"
    ]
    
    for achievement in achievements:
        p = tf9.add_paragraph()
        p.text = achievement
        p.level = 1
        p.font.size = Pt(18)
    
    # Future work
    p_future = tf9.add_paragraph()
    p_future.text = "Future Enhancements"
    p_future.font.size = Pt(22)
    p_future.font.bold = True
    p_future.font.color.rgb = ACCENT_COLOR
    
    future_items = [
        "Support for FBX, STL, and other 3D formats",
        "Blockchain-based verification system",
        "Cloud storage for artist profiles",
        "Batch processing capabilities",
        "Industry standardization efforts"
    ]
    
    for item in future_items:
        p = tf9.add_paragraph()
        p.text = f"• {item}"
        p.level = 1
        p.font.size = Pt(18)
    
    # Save presentation
    output_file = "3D_Model_Authentication_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"\n📊 Presentation Details:")
    print(f"   - Total Slides: 9")
    print(f"   - Format: PowerPoint (.pptx)")
    print(f"   - Slide 1: Title Slide")
    print(f"   - Slide 2: Problem Statement")
    print(f"   - Slide 3: Proposed Solution")
    print(f"   - Slide 4: System Architecture")
    print(f"   - Slide 5: Key Features")
    print(f"   - Slide 6: Technical Implementation")
    print(f"   - Slide 7: Steganographic Algorithm")
    print(f"   - Slide 8: Security Features")
    print(f"   - Slide 9: Conclusion & Future Work")
    print(f"\n✨ Professional design with consistent color scheme and formatting")

if __name__ == "__main__":
    create_presentation()
